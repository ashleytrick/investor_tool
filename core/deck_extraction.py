"""Pitch-deck text extraction + LLM-driven company-profile draft
(Build Session 15).

Two surfaces:

- `extract_text(filename, content_bytes)` -- pure text extraction.
  Tries PDF via pypdf first, then PPTX via python-pptx. Returns
  blocks of text labeled "--- SLIDE N ---" / "--- PAGE N ---" so the
  LLM can localize claims back to a deck position.

- `extract_profile_draft(...)` -- runs the LLM over the extracted
  text and returns a `DeckLLMOutput`. Stub-mode safe (callers pass a
  stub_response when no ANTHROPIC_API_KEY is configured).

This module deliberately does NOT touch SQLite, the workspace
filesystem, or company.yaml. Persistence (if any) is the caller's
problem -- the spec says "do not persist uploaded decks". The web
endpoint reads bytes from the multipart upload, calls these
functions, returns the result, and the bytes drop on the floor.
"""
from __future__ import annotations

import io
import pathlib
from dataclasses import dataclass

from core.llm.client import LLMClient
from schemas.deck_extraction import DeckLLMOutput

# Max chars of deck text to send to the LLM. Pitch decks are short;
# 200k chars is a hard upper bound (Sonnet's context is far larger
# but the prompt only needs the punchy bits). We keep the BEGINNING
# (text[:budget]) because company name + one-liner live on the
# cover slide. Review item #20 fixed an inverted truncation that
# was dropping the start of very-large decks instead of the end.
_LLM_TEXT_BUDGET = 200_000

# Min chars of deck text before we even bother calling the LLM.
# Below this, the deck is image-only or empty -- return warnings and
# let the operator fill in manually.
_MIN_TEXT_FOR_LLM = 200

_PROMPT_PATH = (
    pathlib.Path(__file__).resolve().parent.parent
    / "prompts" / "deck_extraction.txt"
)


@dataclass
class ExtractedText:
    """Plain-text result of parsing the deck file. `text` is what
    we ship to the LLM; `block_count` lets the endpoint surface
    "we found N slides/pages" diagnostics."""
    text: str
    block_count: int
    file_type: str    # "pdf" | "pptx" | "unsupported"
    warnings: list[str]


def extract_text(filename: str, content: bytes) -> ExtractedText:
    """Dispatch by file extension. Legacy .ppt (binary, pre-2007) is
    out of scope -- it requires libreoffice and is unusual today;
    we surface a warning and ask the operator to save as .pptx."""
    lower = (filename or "").lower()
    if lower.endswith(".pdf"):
        return _extract_pdf(content)
    if lower.endswith(".pptx"):
        return _extract_pptx(content)
    if lower.endswith(".ppt"):
        return ExtractedText(
            text="", block_count=0, file_type="unsupported",
            warnings=[
                "legacy .ppt format is not supported; please save the "
                "deck as .pptx (PowerPoint 2007+) and re-upload"
            ],
        )
    return ExtractedText(
        text="", block_count=0, file_type="unsupported",
        warnings=[
            f"unsupported file extension {lower or '(none)'}; "
            "upload a .pdf or .pptx pitch deck"
        ],
    )


def _extract_pdf(content: bytes) -> ExtractedText:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ExtractedText(
            text="", block_count=0, file_type="pdf",
            warnings=[
                "pypdf not installed -- run `uv sync --extra api`"
            ],
        )
    warnings: list[str] = []
    try:
        reader = PdfReader(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001 - pypdf raises diverse types
        return ExtractedText(
            text="", block_count=0, file_type="pdf",
            warnings=[f"could not parse PDF: {exc}"],
        )
    blocks: list[str] = []
    empty_pages = 0
    for i, page in enumerate(reader.pages, 1):
        try:
            page_text = page.extract_text() or ""
        except Exception:  # noqa: BLE001 - per-page failures shouldn't fail the whole deck
            page_text = ""
        page_text = page_text.strip()
        if not page_text:
            empty_pages += 1
            continue
        blocks.append(f"--- PAGE {i} ---\n{page_text}")
    text = "\n\n".join(blocks)
    block_count = len(reader.pages)
    if empty_pages and empty_pages == block_count:
        warnings.append(
            "PDF contains no extractable text -- looks image-only. "
            "Re-export with selectable text, or fill the form manually."
        )
    elif empty_pages > block_count * 0.5:
        warnings.append(
            f"{empty_pages}/{block_count} PDF pages have no "
            f"extractable text (likely image-heavy slides). Review "
            f"extracted fields carefully."
        )
    return ExtractedText(
        text=text, block_count=block_count, file_type="pdf",
        warnings=warnings,
    )


def _extract_pptx(content: bytes) -> ExtractedText:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractedText(
            text="", block_count=0, file_type="pptx",
            warnings=[
                "python-pptx not installed -- run `uv sync --extra api`"
            ],
        )
    warnings: list[str] = []
    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001 - python-pptx raises diverse types
        return ExtractedText(
            text="", block_count=0, file_type="pptx",
            warnings=[f"could not parse PPTX: {exc}"],
        )
    blocks: list[str] = []
    empty_slides = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            # text_frame is the canonical text container in PPTX;
            # tables + grouped shapes need separate handling, but
            # 95% of pitch-deck text lives directly in text_frames.
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                txt = shape.text_frame.text or ""
            except Exception:  # noqa: BLE001
                txt = ""
            txt = txt.strip()
            if txt:
                slide_chunks.append(txt)
        joined = "\n".join(slide_chunks).strip()
        if not joined:
            empty_slides += 1
            continue
        blocks.append(f"--- SLIDE {i} ---\n{joined}")
    text = "\n\n".join(blocks)
    block_count = len(prs.slides)
    if empty_slides and empty_slides == block_count:
        warnings.append(
            "PPTX contains no extractable text -- all slides are "
            "image-only. Add slide titles / body text and re-upload, "
            "or fill the form manually."
        )
    elif empty_slides > block_count * 0.5:
        warnings.append(
            f"{empty_slides}/{block_count} slides have no extractable "
            f"text (likely image-only). Review extracted fields "
            f"carefully."
        )
    return ExtractedText(
        text=text, block_count=block_count, file_type="pptx",
        warnings=warnings,
    )


def extract_profile_draft(
    *, llm: LLMClient, deck_text: str,
    stub_response: dict | None = None,
) -> DeckLLMOutput:
    """Ask the LLM to extract a CompanyProfile draft + per-field
    evidence from the deck text. Returns an empty result + a
    warning when the text is too short to bother extracting."""
    text = (deck_text or "").strip()
    if len(text) < _MIN_TEXT_FOR_LLM:
        return DeckLLMOutput(
            extracted_fields=[],
            warnings=[
                "deck contains too little extractable text to attempt "
                "automatic extraction; please fill the form manually"
            ],
        )
    # Review item #20: keep the BEGINNING. The cover/title slide
    # carries the company name + one-liner; truncating from the
    # start dropped exactly the fields most likely to be required
    # downstream. Traction / ask slides further into the deck are
    # less load-bearing for the extraction (the LLM can still
    # infer "stage" from the first ~50 pages of context).
    truncated = text[:_LLM_TEXT_BUDGET] if len(text) > _LLM_TEXT_BUDGET else text
    prompt = _PROMPT_PATH.read_text(encoding="utf-8").replace(
        "{DECK_TEXT}", truncated,
    )
    return llm.complete_json(
        prompt=prompt,
        schema=DeckLLMOutput,
        stub_response=stub_response,
    )
