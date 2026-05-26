"""Tests for Build Session 15 -- deck-first onboarding.

End-to-end coverage:
- PDF + PPTX bytes are extracted correctly (real bytes generated
  in-test via pypdf/python-pptx so the extractors run for real).
- POST /config/company/extract-from-deck returns the documented
  shape with profile / extracted_fields / missing_required_fields /
  needs_review_fields / warnings / text_preview.
- Endpoint does NOT mutate company.yaml -- the existing PUT path
  remains the only writer.
- Manual PUT /config/company still works for the operator who skips
  the deck upload or hits an extraction warning.
- Unsupported file extensions return a clean warning + empty profile.
"""
from __future__ import annotations

import io
import os
import shutil
from pathlib import Path

import pytest
from fastapi.testclient import TestClient

from tests.conftest import REPO_ROOT


# Long-enough deck text so extract_profile_draft clears the
# _MIN_TEXT_FOR_LLM floor (200 chars) and the stub LLM actually gets
# called. Real pitch decks are well above this floor; the test
# fixtures are intentionally padded to mirror that.
_REALISTIC_DECK_PAGE = (
    "Acme is a B2B compliance API for regulated fintechs. Founded by Jane "
    "Founder, CEO. We turn manual SEC reporting into one POST request. "
    "Currently $200K ARR across four paying design partners. Raising a "
    "seed round to ship self-serve onboarding for the three largest "
    "reporting regimes."
)


# ---------- helpers: build real PDF + PPTX bytes ------------------------

def _make_pdf_bytes(pages: list[str]) -> bytes:
    """Build a minimal in-memory PDF with one page per string. Uses
    pypdf's writer + ReportLab-free path: we render each page via
    pypdf's basic drawing pipeline, which keeps the test free of
    extra deps."""
    from pypdf import PdfWriter
    from pypdf.generic import (
        DecodedStreamObject,
        NameObject,
        TextStringObject,
    )
    writer = PdfWriter()
    for text in pages:
        page = writer.add_blank_page(width=612, height=792)
        # Inject a content stream that draws text. pypdf doesn't
        # expose a high-level "draw text" method on a blank page,
        # so we craft the PostScript-style content stream by hand.
        # This is enough for `page.extract_text()` to find it.
        content = (
            b"BT /F1 12 Tf 72 720 Td ("
            + text.replace("(", "[").replace(")", "]").encode("utf-8")
            + b") Tj ET"
        )
        stream = DecodedStreamObject()
        stream.set_data(content)
        page[NameObject("/Contents")] = stream
        # Minimal font resource so the operator-readable text decoder
        # has something to look up the /F1 reference against.
        page[NameObject("/Resources")] = page.get("/Resources", {})
        # pypdf's add_blank_page already gives /Resources with /Font;
        # if not, force one in here for robustness.
        if "/Font" not in page["/Resources"]:
            from pypdf.generic import DictionaryObject, ArrayObject
            page["/Resources"][NameObject("/Font")] = DictionaryObject({
                NameObject("/F1"): DictionaryObject({
                    NameObject("/Type"): NameObject("/Font"),
                    NameObject("/Subtype"): NameObject("/Type1"),
                    NameObject("/BaseFont"): NameObject("/Helvetica"),
                }),
            })
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _make_pptx_bytes(slides: list[tuple[str, str]]) -> bytes:
    """slides = [(title, body), ...] -> in-memory PPTX bytes."""
    from pptx import Presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # Title Only layout
    for title_text, body_text in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title_text
        # Add a second text box for the body so extract_text picks
        # up more than just the title.
        from pptx.util import Inches
        tb = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4),
        )
        tb.text_frame.text = body_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------- text extraction (pure, no HTTP) -----------------------------

def test_extract_text_pdf_yields_per_page_blocks() -> None:
    from core.deck_extraction import extract_text
    pdf_bytes = _make_pdf_bytes([
        "Acme is a B2B compliance API.",
        "Founded by Jane Founder.",
    ])
    result = extract_text("acme.pdf", pdf_bytes)
    assert result.file_type == "pdf"
    assert result.block_count == 2
    assert "Acme is a B2B compliance API" in result.text
    assert "Jane Founder" in result.text
    # Page tags must survive so the LLM can localize claims back to
    # a deck position.
    assert "--- PAGE 1 ---" in result.text
    assert "--- PAGE 2 ---" in result.text


def test_extract_text_pptx_yields_per_slide_blocks() -> None:
    from core.deck_extraction import extract_text
    pptx_bytes = _make_pptx_bytes([
        ("Acme", "Compliance API for fintechs"),
        ("Why now", "New SEC mandates this quarter"),
    ])
    result = extract_text("acme.pptx", pptx_bytes)
    assert result.file_type == "pptx"
    assert result.block_count == 2
    assert "--- SLIDE 1 ---" in result.text
    assert "--- SLIDE 2 ---" in result.text
    assert "Compliance API for fintechs" in result.text
    assert "New SEC mandates this quarter" in result.text


def test_extract_text_unsupported_extension_warns_cleanly() -> None:
    """Operator uploads a .docx by mistake. Endpoint must NOT crash
    -- return a warning + empty text so the frontend tells the
    operator what went wrong."""
    from core.deck_extraction import extract_text
    result = extract_text("notes.docx", b"\x00\x00")
    assert result.file_type == "unsupported"
    assert result.text == ""
    assert any("unsupported" in w.lower() for w in result.warnings)


def test_extract_text_legacy_ppt_warns_to_save_as_pptx() -> None:
    from core.deck_extraction import extract_text
    result = extract_text("legacy.ppt", b"junk")
    assert result.file_type == "unsupported"
    # Specific .ppt warning so the operator knows the fix (save as
    # .pptx), not just "unsupported".
    assert any("pptx" in w.lower() for w in result.warnings)


def test_extract_text_corrupt_pdf_returns_warning_not_exception() -> None:
    """A garbage PDF body shouldn't crash the endpoint. Return a
    warning + empty text so the operator falls back to manual."""
    from core.deck_extraction import extract_text
    result = extract_text("broken.pdf", b"%PDF-bogus")
    assert result.file_type == "pdf"
    assert result.text == ""
    assert any("could not parse" in w.lower() for w in result.warnings)


def test_extract_profile_draft_skips_llm_when_text_is_sparse() -> None:
    """No point sending an empty / one-line deck to the LLM. The
    builder skips the call and returns a warning so the operator
    knows to fill the form manually."""
    from core.config_loader import load_workspace
    from core.deck_extraction import extract_profile_draft
    from core.llm.client import LLMClient

    ws_src = REPO_ROOT / "clients" / "test_workspace"
    # We don't need a writable workspace -- LLMClient just reads env.
    ws = load_workspace(str(ws_src))
    llm = LLMClient(workspace=ws)
    out = extract_profile_draft(
        llm=llm, deck_text="too short",
        stub_response={"extracted_fields": [], "warnings": []},
    )
    assert out.extracted_fields == []
    assert any("manually" in w.lower() for w in out.warnings)


# ---------- /config/company/extract-from-deck endpoint -------------------

@pytest.fixture
def client(workspace: Path, monkeypatch) -> TestClient:
    """TestClient bound to a fresh workspace + stub LLM."""
    monkeypatch.setenv("API_KEY", "test-api-key")
    monkeypatch.setenv("INVESTOR_WORKSPACE", str(workspace))
    monkeypatch.setenv("CORS_ORIGINS", "")
    monkeypatch.setenv("ANTHROPIC_API_KEY", "")  # stub mode
    import importlib
    import web.api as api_mod
    importlib.reload(api_mod)
    return TestClient(api_mod.app)


def _auth_headers() -> dict[str, str]:
    return {"Authorization": "Bearer test-api-key"}


def test_extract_endpoint_returns_documented_shape_for_pdf(
    client: TestClient, workspace: Path,
) -> None:
    """Happy-path PDF upload returns the response shape the frontend
    is built against -- profile, extracted_fields, missing_required,
    needs_review, warnings, text_preview."""
    pdf_bytes = _make_pdf_bytes([
        _REALISTIC_DECK_PAGE,
        "Traction: 4 design partners. Net revenue retention 128%.",
    ])
    res = client.post(
        "/config/company/extract-from-deck",
        headers=_auth_headers(),
        files={"file": ("acme.pdf", pdf_bytes, "application/pdf")},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    # Shape contract.
    for key in (
        "profile", "extracted_fields", "missing_required_fields",
        "needs_review_fields", "warnings", "source_filename",
        "text_preview",
    ):
        assert key in body
    assert body["source_filename"] == "acme.pdf"
    # Stub LLM populates `name` + `one_liner` + low-conf `problem`.
    assert body["profile"]["name"] == "Stub Co"
    assert "problem" in body["needs_review_fields"], (
        "problem extracted at confidence 0.5 must land in needs_review"
    )
    # founder_email is required for setup but not extracted by the
    # stub -- must surface in missing_required_fields.
    assert "founder_email" in body["missing_required_fields"]
    # Stub-mode warning must propagate.
    assert any("stub" in w.lower() for w in body["warnings"])
    assert "Acme" in body["text_preview"]


def test_extract_endpoint_does_not_mutate_company_yaml(
    client: TestClient, workspace: Path,
) -> None:
    """Hard guarantee from the spec: the extract endpoint must NEVER
    write to company.yaml. Only PUT /config/company does."""
    yaml_path = workspace / "config" / "company.yaml"
    before = yaml_path.read_text(encoding="utf-8")
    pptx_bytes = _make_pptx_bytes([
        ("Acme", _REALISTIC_DECK_PAGE),
        ("Traction", "$200K ARR. 128% net revenue retention. Four "
                     "paying design partners with signed contracts."),
    ])
    res = client.post(
        "/config/company/extract-from-deck",
        headers=_auth_headers(),
        files={"file": (
            "acme.pptx", pptx_bytes,
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
        )},
    )
    assert res.status_code == 200, res.text
    after = yaml_path.read_text(encoding="utf-8")
    # Byte-for-byte unchanged.
    assert after == before


def test_extract_endpoint_then_manual_put_persists(
    client: TestClient, workspace: Path,
) -> None:
    """Full happy-path workflow: upload deck -> review -> PUT. The
    existing PUT path must still work end-to-end after the extract
    response is returned."""
    yaml_path = workspace / "config" / "company.yaml"
    pdf_bytes = _make_pdf_bytes([_REALISTIC_DECK_PAGE])
    ext = client.post(
        "/config/company/extract-from-deck",
        headers=_auth_headers(),
        files={"file": ("acme.pdf", pdf_bytes, "application/pdf")},
    ).json()
    # Operator reviews + edits: fill in the required fields the
    # extraction didn't get to.
    profile = ext["profile"]
    profile["name"] = "Acme Edited"
    profile["one_liner"] = "B2B compliance API."
    profile["founder_name"] = "Jane"
    profile["founder_email"] = "jane@acme.example"
    profile["stage"] = "Seed"
    profile["problem"] = "Manual reporting is slow."
    profile["solution"] = "API for reporting."
    profile["traction"] = "$200K ARR"
    profile["target_sectors"] = ["fintech"]
    profile["scheduling_link"] = "https://cal.example/jane"
    res = client.put(
        "/config/company", headers=_auth_headers(), json=profile,
    )
    assert res.status_code == 200, res.text
    after = yaml_path.read_text(encoding="utf-8")
    assert "Acme Edited" in after
    assert "jane@acme.example" in after


def test_extract_endpoint_unsupported_file_returns_clean_warnings(
    client: TestClient,
) -> None:
    """A .docx or other unsupported file must return 200 with
    warnings, NOT a 4xx -- the frontend should still pre-fill an
    empty form so the operator can continue manually."""
    res = client.post(
        "/config/company/extract-from-deck",
        headers=_auth_headers(),
        files={"file": ("notes.docx", b"not a deck", "application/octet-stream")},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert any("unsupported" in w.lower() for w in body["warnings"])
    # Empty profile + every required field missing.
    assert len(body["missing_required_fields"]) > 0


def test_extract_endpoint_empty_upload_returns_400(
    client: TestClient,
) -> None:
    """A zero-byte upload is a bug, not a normal user case. 400 so
    the frontend surfaces the error rather than rendering an empty
    'extraction succeeded' state."""
    res = client.post(
        "/config/company/extract-from-deck",
        headers=_auth_headers(),
        files={"file": ("empty.pdf", b"", "application/pdf")},
    )
    assert res.status_code == 400, res.text


def test_extract_endpoint_requires_auth() -> None:
    """No Bearer header -> 401. Setup matters for production where
    the wizard pings this with the same auth as the rest of the
    onboarding endpoints."""
    import importlib
    os.environ["API_KEY"] = "test-api-key"
    os.environ["INVESTOR_WORKSPACE"] = str(
        REPO_ROOT / "clients" / "test_workspace"
    )
    import web.api as api_mod
    importlib.reload(api_mod)
    cli = TestClient(api_mod.app)
    res = cli.post(
        "/config/company/extract-from-deck",
        files={"file": ("x.pdf", b"%PDF-1.4\n", "application/pdf")},
    )
    assert res.status_code == 401


def test_openapi_lists_extract_endpoint(client: TestClient) -> None:
    """Lovable regenerates types from /openapi.json; refuse a
    regression that drops the path."""
    spec = client.get("/openapi.json").json()
    assert "/config/company/extract-from-deck" in spec["paths"]
